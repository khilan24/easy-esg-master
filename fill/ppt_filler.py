#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 模板填充：按 JSON 替换 {{占位符}}（与 Word 一致，另增 {{报告日期}}）。
优先 python-pptx（run 级替换 + <a:br/> 换行），无依赖时回退 XML 解包/替换/打包。
"""
import re
import shutil
import zipfile
from pathlib import Path

try:
    from core.utils import find_latest_report_json
except ImportError:
    find_latest_report_json = None

from .word_filler import load_json_report, build_replacements


def _have_pptx():
    try:
        __import__("pptx")
        return True
    except ImportError:
        return False


def _default_template_path():
    p = Path("templates/ESG研报模板.pptx")
    return p if p.exists() else Path("ESG研报模板.pptx")


def _date_part_from_stem(stem):
    if stem.endswith("_报告"):
        return stem[:-3]
    if stem.startswith("报告_"):
        return stem[3:]
    return stem


def _pptx_zip_order(path, temp_dir):
    """打包 pptx 时文件顺序：Content_Types → _rels → docProps → ppt。"""
    rel = path.relative_to(temp_dir)
    p0 = rel.parts[0] if rel.parts else ""
    prio = 0 if p0 == "[Content_Types].xml" else 1 if p0 == "_rels" else 2 if p0 == "docProps" else 3 if p0 == "ppt" else 4
    return (prio, str(rel))


def _norm_val(val):
    if not isinstance(val, str):
        val = str(val) if val is not None else ""
    val = val.strip("\n\r \t")
    return re.sub(r"\n{2,}", "\n", val)


def _val_to_pptx_text(val):
    """返回规范后的替换值（保留 \\n，供段落内换行逻辑使用）。"""
    if val is None:
        return " "
    return _norm_val(val) if isinstance(val, str) else _norm_val(str(val))


def _escape_xml(text):
    """转义 XML 文本/属性中的特殊字符。"""
    if not text:
        return ""
    s = str(text)
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;").replace("'", "&apos;")


def _insert_line_breaks_after_run(para, after_run, lines_after_first):
    """在 after_run 后插入 <a:br/> 和新 run，使「资料来源」等单独成行。"""
    if not lines_after_first:
        return
    try:
        from pptx.oxml import parse_xml
        from copy import deepcopy
    except ImportError:
        return
    p_el = para._element
    r_el = after_run._r
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def make_br():
        return parse_xml(f'<a:br xmlns:a="{ns}"/>')

    def make_run(text):
        rPr = r_el.find(f"{{{ns}}}rPr")
        r_el_str = f'<a:r xmlns:a="{ns}"><a:t>{_escape_xml(text)}</a:t></a:r>'
        run_el = parse_xml(r_el_str)
        if rPr is not None:
            run_el.insert(0, deepcopy(rPr))
        return run_el

    idx = list(p_el).index(r_el)
    for line in lines_after_first:
        p_el.insert(idx + 1, make_br())
        idx += 1
        p_el.insert(idx + 1, make_run(line))
        idx += 1


def _replace_in_paragraph(para, replacements):
    """Run 级占位符替换；若替换值含 \\n 则插入 <a:br/> + 新 run 实现换行。"""
    runs = list(para.runs)
    if not runs:
        return

    for run in runs:
        text = run.text
        if "{{" not in text or "}}" not in text:
            continue
        for ph, val in replacements.items():
            text = text.replace("{{" + ph + "}}", _val_to_pptx_text(val))
        text = re.sub(r"\{\{[^}]*\}\}", " ", text)
        if "\n" in text:
            lines = [s if s.strip() else " " for s in text.split("\n")]
            run.text = lines[0]
            _insert_line_breaks_after_run(para, run, lines[1:])
        else:
            run.text = text
        break

    runs = list(para.runs)
    full = "".join(r.text for r in runs)
    while "{{" in full and "}}" in full:
        start = full.index("{{")
        end = full.find("}}", start) + 2
        if end <= start:
            break
        ph_name = re.sub(r"<[^>]+>", "", full[start + 2 : end - 2]).strip()
        val = replacements.get(ph_name) if ph_name else None
        if val is not None:
            val = _val_to_pptx_text(val)
        else:
            val = " "

        pos = 0
        run_ranges = []
        for r in runs:
            run_ranges.append((r, pos, pos + len(r.text)))
            pos += len(r.text)

        i0 = next(i for i, (_, a, b) in enumerate(run_ranges) if a <= start < b)
        i1 = next(i for i, (_, a, b) in enumerate(run_ranges) if a < end <= b)
        r0, s0, e0 = run_ranges[i0]
        r1, s1, e1 = run_ranges[i1]
        prefix, suffix = full[s0:start], full[end:e1] if i0 == i1 else full[end:e1]
        if "\n" not in val:
            if i0 == i1:
                r0.text = prefix + val + suffix
            else:
                r0.text = prefix + val
                for i in range(i0 + 1, i1):
                    run_ranges[i][0].text = ""
                r1.text = suffix
        else:
            lines = [s if s.strip() else " " for s in val.split("\n")]
            if i0 == i1:
                r0.text = prefix + lines[0]
                rest = lines[1:]
                if rest:
                    last = rest[-1] + suffix
                    _insert_line_breaks_after_run(para, r0, rest[:-1] + [last])
                elif suffix:
                    _insert_line_breaks_after_run(para, r0, [suffix])
            else:
                r0.text = prefix + lines[0]
                for i in range(i0 + 1, i1):
                    run_ranges[i][0].text = ""
                r1.text = suffix
                if lines[1:]:
                    _insert_line_breaks_after_run(para, r0, lines[1:])
        runs = list(para.runs)
        full = "".join(r.text for r in runs)


def _process_shape(shape, replacements):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            _replace_in_paragraph(para, replacements)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    _replace_in_paragraph(para, replacements)
    if getattr(shape, "shapes", None):
        for s in shape.shapes:
            _process_shape(s, replacements)


def _slide_xml_has_no_content(xml_content):
    """幻灯片 XML 中所有 a:t 文本为空或仅空白则视为无内容。"""
    text = "".join(re.findall(r"<a:t>([^<]*)</a:t>", xml_content))
    return not text.strip()


def _delete_empty_slides_from_pptx(pptx_path):
    """从已保存的 pptx 中删除无内容幻灯片（解包→删幻灯→重打包）。"""
    pptx_path = Path(pptx_path)
    temp_dir = pptx_path.parent / ("_temp_ppt_trim_" + pptx_path.stem)
    out_temp = None
    if temp_dir.exists():
        shutil.rmtree(temp_dir)
    temp_dir.mkdir(parents=True)
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            zf.extractall(temp_dir)
        slides_dir = temp_dir / "ppt" / "slides"
        to_delete = []
        for xml_file in sorted(slides_dir.glob("slide*.xml"), key=lambda p: int(re.search(r"(\d+)", p.name).group(1))):
            try:
                content = xml_file.read_text(encoding="utf-8")
                if _slide_xml_has_no_content(content):
                    to_delete.append(xml_file)
            except Exception:
                continue
        _remove_slides(temp_dir, to_delete)
        out_temp = pptx_path.parent / (pptx_path.stem + "_out.pptx")
        with zipfile.ZipFile(out_temp, "w", zipfile.ZIP_DEFLATED) as zf:
            for f in sorted((x for x in temp_dir.rglob("*") if x.is_file()), key=lambda x: _pptx_zip_order(x, temp_dir)):
                zf.write(f, f.relative_to(temp_dir))
        out_temp.replace(pptx_path)
    finally:
        if temp_dir.exists():
            shutil.rmtree(temp_dir)
        if out_temp is not None and out_temp.exists():
            try:
                out_temp.unlink()
            except Exception:
                pass


def _fill_via_pptx(template_path, output_path, replacements):
    """使用 python-pptx 做 run 级替换并保存，删除无内容幻灯。"""
    from pptx import Presentation
    prs = Presentation(str(template_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            _process_shape(shape, replacements)
    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            _process_shape(shape, replacements)
    for master in prs.slide_masters:
        for shape in master.shapes:
            _process_shape(shape, replacements)
    prs.save(str(output_path))
    _delete_empty_slides_from_pptx(output_path)
    return True


def _build_ppt_replacements(report_data, max_news_per_section=8):
    replacements = build_replacements(report_data, max_news_per_section)
    gen_time = report_data.get("report_metadata", {}).get("generation_time", "")
    if gen_time:
        replacements["报告日期"] = gen_time.split()[0] if " " in gen_time else gen_time
    return replacements


# ---------------------------------------------------------------------------
# XML 回退路径（无 python-pptx 时）
# ---------------------------------------------------------------------------

def _to_pptx_text(text):
    """DrawingML 安全文本：控制字符替换、转义、换行 → </a:t><a:br/><a:t>。"""
    if not isinstance(text, str):
        text = str(text)
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", " ", text)
    text = _escape_xml(text).strip("\n\r \t")
    if not text:
        return " "
    text = re.sub(r"\n{2,}", "\n", text)
    return text.replace("\n", "</a:t><a:br/><a:t>")


def _replace_in_xml(xml_content, placeholder, replacement):
    """占位符替换：先字面 {{name}}，否则扫描 {{...}} 整段替换。"""
    repl = _to_pptx_text(replacement)
    ph_text = "{{" + placeholder + "}}"
    if ph_text in xml_content:
        return xml_content.replace(ph_text, repl), True
    i = 0
    while i < len(xml_content) - 1:
        if xml_content[i : i + 2] != "{{":
            i += 1
            continue
        depth, pos, end = 0, i + 2, -1
        while pos < len(xml_content) - 1:
            if xml_content[pos : pos + 2] == "}}":
                if depth == 0:
                    end = pos + 2
                    break
                depth -= 1
            elif xml_content[pos : pos + 2] == "{{":
                depth += 1
            pos += 1
        if end > 0:
            inner = re.sub(r"<[^>]+>", "", xml_content[i + 2 : end - 2])
            if placeholder in inner or inner.strip() == placeholder:
                return xml_content[:i] + repl + xml_content[end:], True
        i = end + 2 if end > 0 else i + 1
    return xml_content, False


def _clean_xml_newlines(xml_content):
    xml_content = re.sub(r"(</a:t><a:br/><a:t>){2,}", "</a:t><a:br/><a:t>", xml_content)
    xml_content = re.sub(r"<a:t></a:t><a:br/><a:t></a:t>", "<a:t> </a:t>", xml_content)
    xml_content = re.sub(r"^(</a:t><a:br/><a:t>)+", "", xml_content)
    xml_content = re.sub(r"(</a:t><a:br/><a:t>)+$", "", xml_content)
    xml_content = re.sub(r"<a:t></a:t>", "<a:t> </a:t>", xml_content)
    return xml_content


def _clear_remaining_placeholders(xml_content, used):
    i = 0
    while i < len(xml_content) - 3:
        if xml_content[i : i + 2] != "{{":
            i += 1
            continue
        end = xml_content.find("}}", i + 2)
        if end == -1:
            i += 1
            continue
        name = re.sub(r"<[^>]+>", "", xml_content[i + 2 : end]).strip()
        if name and name not in used and xml_content[i : end + 2] == "{{" + name + "}}":
            xml_content = xml_content[:i] + " " + xml_content[end + 2 :]
            i += 1
            continue
        i = end + 2
    return xml_content


def _extract_placeholders(xml_content):
    out = set()
    i = 0
    while i < len(xml_content) - 1:
        if xml_content[i : i + 2] != "{{":
            i += 1
            continue
        pos = xml_content.find("}}", i + 2)
        if pos == -1:
            i += 1
            continue
        inner = re.sub(r"<[^>]+>", "", xml_content[i + 2 : pos]).strip()
        if inner:
            out.add(inner)
        i = pos + 2
    return out


def _is_empty_news_slide(xml_content, used_placeholders):
    phs = _extract_placeholders(xml_content)
    if not phs or any(p in used_placeholders for p in phs):
        return False
    return all(re.match(r"^(环境|社会|治理)新闻(标题|内容)\d+$", p) for p in phs)


def _remove_slides(temp_dir, slide_files):
    if not slide_files:
        return
    slides_dir = temp_dir / "ppt" / "slides"
    pres_xml = temp_dir / "ppt" / "presentation.xml"
    pres_rels = temp_dir / "ppt" / "_rels" / "presentation.xml.rels"
    content_types = temp_dir / "[Content_Types].xml"
    app_xml = temp_dir / "docProps" / "app.xml"
    try:
        rels_content = pres_rels.read_text(encoding="utf-8")
    except Exception:
        return
    rids = []
    for f in slide_files:
        m = re.search(
            rf'<Relationship\s+Id="(rId\d+)"[^>]*Target="[^"]*{re.escape(f.name)}"',
            rels_content,
        )
        if m:
            rids.append((m.group(1), f))
    if not rids:
        return
    pres = pres_xml.read_text(encoding="utf-8")
    rels = rels_content
    for r_id, f in rids:
        pres = re.sub(rf'<p:sldId\s+[^>]*r:id="{re.escape(r_id)}"[^>]*/>\s*', "", pres)
        rels = re.sub(rf'<Relationship\s+Id="{re.escape(r_id)}"[^>]*/>\s*', "", rels, flags=re.DOTALL)
        f.unlink(missing_ok=True)
        (slides_dir / "_rels" / (f.name + ".rels")).unlink(missing_ok=True)
    pres_xml.write_text(pres, encoding="utf-8")
    pres_rels.write_text(rels, encoding="utf-8")
    if content_types.exists():
        ct = content_types.read_text(encoding="utf-8")
        for _, f in rids:
            ct = re.sub(rf'<Override\s+PartName="/ppt/slides/{re.escape(f.name)}"[^/]*/>\s*', "", ct)
        content_types.write_text(ct, encoding="utf-8")
    if app_xml.exists():
        try:
            app = app_xml.read_text(encoding="utf-8")
            n = len(rids)
            app = re.sub(r"<Slides>(\d+)</Slides>", lambda m: f"<Slides>{max(0, int(m.group(1)) - n)}</Slides>", app)
            app_xml.write_text(app, encoding="utf-8")
        except Exception:
            pass


def _fill_via_xml(template_path, output_path, replacements):
    """解包 → 替换占位符与清理 → 删空幻灯 → 打包。"""
    used = set(replacements.keys())
    temp_dir = Path("temp_ppt_unpacked")
    if temp_dir.exists():
        shutil.rmtree(temp_dir)
    temp_dir.mkdir(parents=True)
    with zipfile.ZipFile(template_path, "r") as zf:
        zf.extractall(temp_dir)
    slides_dir = temp_dir / "ppt" / "slides"
    to_delete = []
    for xml_file in temp_dir.rglob("*.xml"):
        s = str(xml_file.relative_to(temp_dir)).replace("\\", "/")
        if not (s.startswith("ppt/slides/slide") or s.startswith("ppt/slideLayouts/")) or not s.endswith(".xml"):
            continue
        try:
            content = xml_file.read_text(encoding="utf-8")
        except Exception:
            continue
        orig = content
        changed = False
        for ph, val in replacements.items():
            content, ok = _replace_in_xml(content, ph, val)
            if ok:
                changed = True
        content = _clean_xml_newlines(content)
        is_content_slide = slides_dir.exists() and slides_dir in xml_file.parents and xml_file.name.startswith("slide") and xml_file.name.endswith(".xml")
        if is_content_slide and not changed and _is_empty_news_slide(content, used):
            to_delete.append(xml_file)
        content = _clear_remaining_placeholders(content, used)
        if content != orig:
            xml_file.write_text(content, encoding="utf-8")
    _remove_slides(temp_dir, to_delete)
    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for f in sorted((x for x in temp_dir.rglob("*") if x.is_file()), key=lambda x: _pptx_zip_order(x, temp_dir)):
            zf.write(f, f.relative_to(temp_dir))
    shutil.rmtree(temp_dir)
    return True


def fill_ppt_template(json_path=None, template_path=None, output_path=None):
    """
    填充 PPT 模板。
    json_path: 报告 JSON，None 时自动查找最新 *_报告.json
    template_path: 模板路径，None 时用 templates/ESG研报模板.pptx 或根目录
    output_path: 输出路径，None 时与 JSON 同目录、{日期}_最终版.pptx
    返回 (success: bool, output_path: Path)
    """
    template_path = Path(template_path) if template_path else _default_template_path()
    output_dir = Path("output")
    output_dir.mkdir(exist_ok=True)
    if json_path is None:
        json_path = find_latest_report_json(output_dir) if find_latest_report_json else None
        if not json_path:
            json_files = list(output_dir.glob("**/*_报告.json")) or list(Path(".").glob("**/*_报告.json"))
            json_path = max(json_files, key=lambda p: p.stat().st_mtime) if json_files else None
        if not json_path or not json_path.exists():
            print("错误：未找到 JSON 报告文件")
            return False, None
    json_path = Path(json_path)
    if output_path is None:
        output_path = json_path.parent / f"{_date_part_from_stem(json_path.stem)}_最终版.pptx"
    else:
        output_path = Path(output_path)
    if not template_path.exists():
        print(f"错误：PPT 模板不存在: {template_path}")
        return False, None
    if not json_path.exists():
        print(f"错误：JSON 文件不存在: {json_path}")
        return False, None
    report_data = load_json_report(json_path)
    replacements = _build_ppt_replacements(report_data, max_news_per_section=8)
    if _have_pptx():
        try:
            _fill_via_pptx(template_path, output_path, replacements)
            return True, output_path
        except Exception as e:
            print(f"[PPT] python-pptx 填充失败，回退 XML: {e}")
    _fill_via_xml(template_path, output_path, replacements)
    return True, output_path
